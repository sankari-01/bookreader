from flask import Flask, render_template, request, redirect, send_from_directory, jsonify
from flask_cors import CORS
import os
from datetime import datetime
import PyPDF2
from pptx import Presentation
from docx import Document
import pytesseract
from PIL import Image
import mammoth
import re
try:
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup
    import mobi
    import html2text
except ImportError:
    pass

from utils.summarize import summarize_text
from utils.qa import answer_question
from utils.speech import text_to_speech
from utils.dictionary import get_meaning
from utils.logger import log_error
from utils.transcribe import transcribe_video
from deep_translator import GoogleTranslator
import mysql.connector
import ocr
import time
import random
from langdetect import detect
import hashlib
import json
import threading

app = Flask(__name__)
CORS(app) # Enable CORS for all routes

def get_db_connection():
    try:
        return mysql.connector.connect(
            host="localhost",
            user="root",
            password="",
            database="ai_books_db",
            connect_timeout=3 # Don't hang the app if DB is down
        )
    except Exception as e:
        # We log it but return None so callers can decide fallback
        log_error(f"Database Connection Failed: {e}")
        return None

UPLOAD_FOLDER = "uploads"
TEXT_CACHE_DIR = "text_cache"
TRANSLATION_CACHE_DIR = "translation_cache"
SUMMARY_CACHE_DIR = "summary_cache"

for folder in [UPLOAD_FOLDER, TEXT_CACHE_DIR, TRANSLATION_CACHE_DIR, SUMMARY_CACHE_DIR, os.path.join("static", "extracted_images")]:
    if not os.path.exists(folder):
        os.makedirs(folder)

def get_file_hash(path):
    """Generates a hash for a file based on its path and modification time."""
    if not os.path.exists(path):
        return None
    mtime = os.path.getmtime(path)
    # Using simple path + mtime hash for speed and to detect file changes
    return hashlib.md5(f"{path}_{mtime}".encode()).hexdigest()

def get_text_cache_path(file_path):
    h = get_file_hash(file_path)
    return os.path.join(TEXT_CACHE_DIR, f"{h}.txt") if h else None

def get_translation_cache_path(file_path, lang):
    h = get_file_hash(file_path)
    return os.path.join(TRANSLATION_CACHE_DIR, f"{h}_{lang}.json") if h else None

# Set Tesseract path for Windows
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

def robust_translate(content, translator, source_lang='auto'):
    """Translates content while preserving structure and handling errors gracefully per line."""
    if not content or not content.strip():
        return content
    
    # Split into paragraphs to preserve structure
    lines = content.split('\n')
    translated_lines = []
    
    for line in lines:
        if not line.strip():
            translated_lines.append(line)
            continue
        
        # Inner helper for a single translation with retries and exponential backoff
        def translate_with_retry(text, retries=5):
            for i in range(retries):
                try:
                    # Jittered delay to stay under Google's radar
                    time.sleep(random.uniform(0.15, 0.4) * (i + 1))
                    result = translator.translate(text)
                    if result and result.strip():
                        return result
                    return text
                except Exception as e:
                    log_error(f"Translation attempt {i+1} failed for text segment: {e}")
                    if i == retries - 1:
                        return text # Final fallback: return original
            return text

        # Handle long lines using sentence-aware chunking
        if len(line) > 4000:
            sub_chunks = re.split(r'(?<=[.!?])\s+', line)
            current_batch = ""
            translated_batch = []
            
            for sc in sub_chunks:
                if len(current_batch) + len(sc) < 4000:
                    current_batch += ( " " if current_batch else "" ) + sc
                else:
                    if current_batch.strip():
                        translated_batch.append(translate_with_retry(current_batch.strip()))
                    current_batch = sc
            
            if current_batch.strip():
                translated_batch.append(translate_with_retry(current_batch.strip()))
            
            translated_lines.append(" ".join(translated_batch))
        else:
            translated_lines.append(translate_with_retry(line))
    
    return "\n".join(translated_lines)

def translate_with_markers(text, lang, source_lang='auto'):
    """Translates text while preserving --- Page X --- and --- Slide X --- markers."""
    if not text:
        return text
    
    marker_regex = r'(\s*--- (?:Page|Slide) \d+ ---\s*)'
    parts = re.split(marker_regex, text)
    
    try:
        # Use explicit source_lang for much higher reliability than 'auto'
        translator = GoogleTranslator(source=source_lang, target=lang)
    except Exception as e:
        log_error(f"Failed to initialize translator: {e}")
        return text

    translated_parts = []
    for i, part in enumerate(parts):
        if re.search(r'--- (?:Page|Slide) \d+ ---', part):
            translated_parts.append(part)
        else:
            # Progress logging for large books
            if i > 0 and i % 5 == 0:
                log_error(f"Translating chunk {i}/{len(parts)} (Target: {lang})")
            translated_parts.append(robust_translate(part, translator, source_lang))
                
    return "".join(translated_parts)

def generate_html_preview(path, filename, lang='en'):
    html = ""
    lower_name = filename.lower()
    
    try:
        translator = GoogleTranslator(source='auto', target=lang) if lang != 'en' else None
    except:
        translator = None

    def translate_if_needed(content):
        if not translator or not content.strip():
            return content
        return robust_translate(content, translator)

    def translate_html_preserving_tags(html_content):
        # We NO LONGER translate the background HTML to ensure "Original View" stays original
        # and to avoid double-translation issues with the overlay.
        return html_content

    try:
        if lower_name.endswith((".docx", ".doc")):
            if lower_name.endswith(".docx"):
                with open(path, "rb") as docx_file:
                    result = mammoth.convert_to_html(docx_file)
                    html = f'<div class="docx-preview">{translate_html_preserving_tags(result.value)}</div>'
            else:
                # Basic DOC support (extract text and wrap in pre)
                try:
                    import subprocess
                    # If antiword is installed, we could use it. Otherwise, fallback to text.
                    with open(path, "r", encoding="utf-8", errors="ignore") as f:
                        html = f'<pre class="doc-preview">{translate_if_needed(f.read())}</pre>'
                except:
                    html = '<div class="error-preview">Detailed preview not available for legacy .doc files.</div>'
        
        elif lower_name.endswith(".pptx"):
            prs = Presentation(path)
            html = '<div class="pptx-preview">'
            for i, slide in enumerate(prs.slides, start=1):
                slide_html = f'<div class="slide-box"><h4>{translate_if_needed("Slide")} {i}</h4><ul>'
                has_content = False
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        translated_text = translate_if_needed(shape.text.strip())
                        slide_html += f"<li>{translated_text}</li>"
                        has_content = True
                
                if not has_content:
                    slide_html += f"<li><i>{translate_if_needed('Empty Slide')}</i></li>"
                
                slide_html += '</ul></div>'
                html += slide_html
            html += '</div>'
            
        elif lower_name.endswith(".txt"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
                if lang != 'en':
                    content = translate_if_needed(content)
                html = f'<pre class="txt-preview" style="white-space: pre-wrap; font-family: inherit;">{content}</pre>'

        elif lower_name.endswith(".epub"):
            book = epub.read_epub(path)
            html = '<div class="epub-preview">'
            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                content = item.get_content().decode("utf-8")
                soup = BeautifulSoup(content, 'html.parser')
                # Remove scripts and styles
                for s in soup(['script', 'style']): s.decompose()
                html += f'<div class="epub-section">{translate_html_preserving_tags(str(soup.body if soup.body else soup))}</div>'
            html += '</div>'

        elif lower_name.endswith(".mobi"):
            try:
                temp_dir = os.path.join(os.getcwd(), 'tmp', f"mobi_{int(time.time())}")
                os.makedirs(temp_dir, exist_ok=True)
                mobi.extract(path, temp_dir)
                # Find the HTML file in extracted folder
                for root, dirs, files in os.walk(temp_dir):
                    for f in files:
                        if f.endswith(('.html', '.htm')):
                            with open(os.path.join(root, f), 'r', encoding='utf-8', errors='ignore') as mf:
                                m_html = mf.read()
                                html = f'<div class="mobi-preview">{translate_html_preserving_tags(m_html)}</div>'
                                break
            except Exception as me:
                html = f'<div class="error-preview">Mobi error: {str(me)}</div>'

        elif lower_name.endswith(".html"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                html = f'<div class="html-preview">{translate_html_preserving_tags(f.read())}</div>'
        
        else:
            # Universal Fallback: Try reading as text
            try:
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
                    if content and len(content) > 10: # Only if it looks like text
                        html = f'<pre class="fallback-text-preview" style="white-space: pre-wrap;">{translate_if_needed(content)}</pre>'
                    else:
                        html = f'<div class="error-preview">Unsupported file type: {filename}</div>'
            except:
                html = f'<div class="error-preview">Unsupported file type: {filename}</div>'
                
    except Exception as e:
        html = f'<div class="error-preview">Preview error: {str(e)}</div>'
    
    return html

def extract_text_from_file(path, filename):
    log_error(f"Starting extraction for: {filename}")
    
    # Check cache first
    cache_path = get_text_cache_path(path)
    if cache_path and os.path.exists(cache_path):
        log_error(f"Using cached text for: {filename}")
        with open(cache_path, "r", encoding="utf-8") as f:
            cached_data = json.load(f)
            return cached_data['text'], cached_data['is_image'], cached_data['is_video']

    text = ""
    is_image = False
    is_video = False
    lower_name = filename.lower()

    try:
        if lower_name.endswith(".pdf"):
            text = ocr.extract_text(path)

        elif lower_name.endswith((".pptx", ".ppt")):
            if lower_name.endswith(".ppt"):
                from utils.office_to_pdf import convert_to_pdf
                pdf_path = convert_to_pdf(path)
                if pdf_path and os.path.exists(pdf_path):
                    text = ocr.extract_text(pdf_path)
                else:
                    text = "Failed to convert PPT to PDF for text extraction."
            else:
                prs = Presentation(path)
                for slide_no, slide in enumerate(prs.slides, start=1):
                    text += f"\n--- Slide {slide_no} ---\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text += shape.text + "\n"
                    text += "\n"

        elif lower_name.endswith((".docx", ".doc")):
            # Fast native extraction instead of slow COM PDF conversion
            if lower_name.endswith(".docx"):
                with open(path, "rb") as docx_file:
                    raw_text = mammoth.extract_raw_text(docx_file).value
                # Chunk long Word docs into "pages" (approx 2500 chars per page) for consistent pagination and TTS
                chunk_size = 2500
                chunks = [raw_text[i:i+chunk_size] for i in range(0, max(1, len(raw_text)), chunk_size)]
                text = ""
                for i, chunk in enumerate(chunks, 1):
                    text += f"\n--- Page {i} ---\n" + chunk.strip() + "\n"
            else:
                text = "Failed to extract text from generic .doc. Please convert to .docx for optimal performance."

        elif lower_name.endswith((".epub", ".mobi")):
            if lower_name.endswith(".epub"):
                book = ebooklib.epub.read_epub(path)
                item_no = 1
                for item in book.get_items():
                    if item.get_type() == ebooklib.ITEM_DOCUMENT:
                        soup = BeautifulSoup(item.get_content(), 'html.parser')
                        text += f"\n--- Page {item_no} ---\n" + soup.get_text() + "\n"
                        item_no += 1
            else:
                # MOBI support
                temp_dir, out_file = mobi.extract(path)
                with open(out_file, 'r', encoding='utf-8', errors='ignore') as f:
                    html_content = f.read()
                    soup = BeautifulSoup(html_content, 'html.parser')
                    full_text = soup.get_text()
                    # Chunk MOBI into "pages" for consistent translation UI
                    chunks = [full_text[i:i+3000] for i in range(0, len(full_text), 3000)]
                    for i, chunk in enumerate(chunks, 1):
                        text += f"\n--- Page {i} ---\n" + chunk + "\n"

        elif lower_name.endswith((".html", ".htm", ".md")):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
                if lower_name.endswith(".md"):
                    text = content # Already text
                else:
                    soup = BeautifulSoup(content, 'html.parser')
                    text = soup.get_text()
                
                # Segment long HTML/MD into pages for translation UI
                chunks = [text[i:i+3000] for i in range(0, len(text), 3000)]
                text = ""
                for i, chunk in enumerate(chunks, 1):
                    text += f"\n--- Page {i} ---\n" + chunk + "\n"

        elif lower_name.endswith((".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp")):
            is_image = True
            image = Image.open(path)
            extracted_text = pytesseract.image_to_string(image).strip()
            text = extracted_text if extracted_text else "No readable text found in this image."

        elif lower_name.endswith((".mp4", ".avi", ".mov", ".mkv", ".webm")):
            is_video = True
            log_error(f"Transcribing video: {filename}")
            text = transcribe_video(path)

        else:
            # Emergency fallback: Try reading as text even if extension is unknown
            try:
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
                    log_error(f"Fallback text extraction for unknown type: {filename}")
            except:
                text = "Preview not supported for this file type."
                log_error(f"Unsupported file type: {filename}")

        if text:
            # 1. Protect structural markers by ensuring they are isolated with double newlines
            # This prevents the subsequent word-wrap fix from collapsing them into the text.
            text = re.sub(r'\s*(--- (?:Page|Slide) \d+ ---)\s*', r'\n\n\1\n\n', text)
            
            # 2. Normalize multiple newlines into double newlines (standard paragraph separation)
            text = re.sub(r'\n\s*\n', '\n\n', text)
            
            # 3. PDF Word-Wrap Fix: Replace single newlines with spaces ONLY if they are NOT
            # part of a marker or a double newline.
            text = re.sub(r'(?<!\n)\n(?!\n)', ' ', text)
            
            # 4. Final Cleanup: Compress multiple spaces and ensure markers are the only thing on their lines
            text = re.sub(r'[ \t]+', ' ', text).strip()
            # Ensure markers are indeed on their own lines even after space compression
            text = re.sub(r' ?(--- (?:Page|Slide) \d+ ---) ?', r'\n\n\1\n\n', text)
            # Final collapse of redundant newlines
            text = re.sub(r'\n{3,}', '\n\n', text).strip()

        # Save to cache if successful
        if text and cache_path:
            with open(cache_path, "w", encoding="utf-8") as f:
                json.dump({"text": text, "is_image": is_image, "is_video": is_video}, f)

        log_error(f"Extraction successful for: {filename} ({len(text)} chars)")
    except Exception as e:
        error_msg = f"Unable to read this file. Error: {str(e)}"
        log_error(f"File: {filename} - {error_msg}")
        text = error_msg

    return text, is_image, is_video

def get_page_count(path, filename):
    lower_name = filename.lower()
    try:
        if lower_name.endswith(".pdf"):
            reader = PyPDF2.PdfReader(path)
            return str(len(reader.pages))
        elif lower_name.endswith(".pptx"):
            prs = Presentation(path)
            return str(len(prs.slides))
        elif lower_name.endswith(".docx"):
            # Instant page count calculation based on chunk length
            try:
                with open(path, "rb") as docx_file:
                    raw_text = mammoth.extract_raw_text(docx_file).value
                    chunk_size = 2500
                    return str(max(1, (len(raw_text) + chunk_size - 1) // chunk_size))
            except Exception:
                return "1"
        elif lower_name.endswith(".doc"):
            return "N/A"
        elif lower_name.endswith((".png", ".jpg", ".jpeg", ".gif", ".webp")):
            return "1"
        elif lower_name.endswith((".txt", ".mp4", ".avi", ".mkv")):
            return "N/A"
        else:
            return "-"
    except Exception:
        return "-"


@app.route("/api/health")
def health():
    return jsonify({"status": "ok"})


@app.route("/api/upload", methods=["POST"])
def upload():
    file = request.files.get("file")

    if not file or file.filename.strip() == "":
        return jsonify({"error": "No file selected"}), 400

    filename = file.filename
    path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(path)

    upload_time = datetime.now()
    pages_count = get_page_count(path, filename)

    db_error = False
    try:
        conn = get_db_connection()
        if conn:
            cursor = conn.cursor()
            cursor.execute("SELECT id FROM uploaded_books WHERE filename = %s", (filename,))
            if cursor.fetchone():
                cursor.execute("UPDATE uploaded_books SET upload_time = %s, pages = %s WHERE filename = %s", (upload_time, pages_count, filename))
            else:
                cursor.execute("INSERT INTO uploaded_books (filename, upload_time, pages) VALUES (%s, %s, %s)", (filename, upload_time, pages_count))
            conn.commit()
            conn.close()
        else:
            db_error = True
    except Exception as e:
        log_error(f"DB Error on upload: {e}")
        db_error = True

    # Even if DB fails, we return success because the file was saved to filesystem at line 437
    return jsonify({"success": True, "filename": filename, "db_error": db_error})

    return jsonify({"success": True, "filename": filename})


@app.route("/api/debug")
def debug_env():
    import os
    return jsonify({
        "cwd": os.getcwd(),
        "upload_folder": UPLOAD_FOLDER,
        "abs_upload_folder": os.path.abspath(UPLOAD_FOLDER),
        "exists": os.path.exists(UPLOAD_FOLDER),
        "listdir": os.listdir(UPLOAD_FOLDER) if os.path.exists(UPLOAD_FOLDER) else "N/A"
    })

@app.route("/api/files")
def files():
    file_info = []
    db_error = False

    try:
        conn = get_db_connection()
        cursor = conn.cursor(dictionary=True)
        cursor.execute("SELECT id, filename, upload_time, pages, last_opened, last_closed FROM uploaded_books ORDER BY upload_time DESC")
        
        for row in cursor.fetchall():
            path = os.path.join(UPLOAD_FOLDER, row['filename'])
            if os.path.exists(path):
                dt = row['upload_time']
                
                date_str = dt.strftime("%d-%m-%Y") if dt else "-"
                day_str = dt.strftime("%A") if dt else "-"
                time_str = dt.strftime("%I:%M %p") if dt else "-"
                
                last_opened = row['last_opened']
                last_closed = row['last_closed']
                
                file_info.append({
                    "id": row['id'],
                    "filename": row['filename'],
                    "date": date_str,
                    "day": day_str,
                    "time": time_str,
                    "pages": row['pages'] or "-",
                    "last_opened": last_opened.strftime("%d-%m-%Y %I:%M %p") if last_opened else "Never",
                    "last_closed": last_closed.strftime("%d-%m-%Y %I:%M %p") if last_closed else "Never"
                })
        conn.close()
    except Exception as e:
        log_error(f"DB Error on fetch files: {e}")
        db_error = True

    # Always sync with filesystem to ensure no files are missed (even if DB is incomplete)
    existing_filenames = {f['filename'] for f in file_info}
    for filename in os.listdir(UPLOAD_FOLDER):
        # Ignore temporary OS/Word lock files
        if filename.startswith('.') or filename.startswith('~$'):
            continue
            
        path = os.path.join(UPLOAD_FOLDER, filename)
        if os.path.isfile(path) and filename not in existing_filenames:
            try:
                mtime = os.path.getmtime(path)
                dt = datetime.fromtimestamp(mtime)
            except Exception:
                dt = datetime.now()
            
            file_info.append({
                "id": "FS",
                "filename": filename,
                "date": dt.strftime("%d-%m-%Y"),
                "day": dt.strftime("%A"),
                "time": dt.strftime("%I:%M %p"),
                "pages": get_page_count(path, filename),
                "last_opened": "Never",
                "last_closed": "Never"
            })
    
    # Sort everything by modification time (most recent first)
    file_info.sort(key=lambda x: os.path.getmtime(os.path.join(UPLOAD_FOLDER, x['filename'])) if os.path.exists(os.path.join(UPLOAD_FOLDER, x['filename'])) else 0, reverse=True)

    return jsonify({"files": file_info, "db_error": db_error})

@app.route("/api/book/timestamp", methods=["POST"])
def update_book_timestamp():
    filename = request.form.get("filename", "").strip()
    ts_type = request.form.get("type", "").strip() # 'opened' or 'closed'
    
    if not filename or ts_type not in ['opened', 'closed']:
        return jsonify({"error": "Invalid request"}), 400
        
    try:
        conn = get_db_connection()
        if not conn:
            return jsonify({"success": False, "error": "Database unavailable"}), 503
            
        cursor = conn.cursor()
        now = datetime.now()
        
        column = "last_opened" if ts_type == "opened" else "last_closed"
        cursor.execute(f"UPDATE uploaded_books SET {column} = %s WHERE filename = %s", (now, filename))
        
        conn.commit()
        conn.close()
        return jsonify({"success": True, "timestamp": now.strftime("%Y-%m-%d %H:%M:%S")})
    except Exception as e:
        log_error(f"Error updating timestamp for {filename}: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/delete/<path:filename>", methods=["POST"])
def delete_book(filename):
    path = os.path.join(UPLOAD_FOLDER, filename)
    
    # Try deleting from database first regardless of file existence
    try:
        conn = get_db_connection()
        if conn:
            cursor = conn.cursor()
            cursor.execute("DELETE FROM uploaded_books WHERE filename = %s", (filename,))
            conn.commit()
            conn.close()
    except Exception as e:
        log_error(f"DB Error on delete: {e}")
        # We don't return 500 here since the file might just be on the filesystem
        # but we log the error.
        
    # Special handling for Docx deletion: attempt to kill any stuck Word processes
    # that might still have a lock (especially after a failed conversion)
    if filename.lower().endswith(".docx"):
        try:
            import subprocess
            subprocess.run(["taskkill", "/F", "/IM", "WINWORD.EXE", "/T"], capture_output=True)
            time.sleep(0.5) # Allow OS to release handles
        except:
            pass
            
    # Then delete the actual file and its associated caches
    try:
        # 1. Delete associated caches
        base, _ = os.path.splitext(path)
        pdf_path = base + ".pdf"
        if os.path.exists(pdf_path):
            try: os.remove(pdf_path)
            except: pass
            
        text_cache = get_text_cache_path(path)
        if text_cache and os.path.exists(text_cache):
            try: os.remove(text_cache)
            except: pass
            
        # Delete all translation caches for this file
        h = get_file_hash(path)
        if h:
            for f in os.listdir(TRANSLATION_CACHE_DIR):
                if f.startswith(h):
                    try: os.remove(os.path.join(TRANSLATION_CACHE_DIR, f))
                    except: pass

        # 2. Delete the main file with a retry loop (handle transient locks)
        if os.path.exists(path):
            success = False
            for attempt in range(3):
                try:
                    os.remove(path)
                    success = True
                    break
                except Exception as e:
                    log_error(f"Delete attempt {attempt+1} failed for {filename}: {e}")
                    time.sleep(0.5) # Wait for other processes to release the file
            
            if not success:
                 return jsonify({"error": "Failed to delete file. It may be in use by another process."}), 500
    except Exception as e:
        log_error(f"General deletion error for {filename}: {e}")
        return jsonify({"error": "Failed to delete file"}), 500
            
    return jsonify({"success": True})


# --- BACKGROUND PRE-TRANSLATION ENGINE ---
# A curated list of supported languages for pre-calculation
PRE_TRANSLATE_LANGS = ['ta', 'hi', 'fr', 'es', 'de', 'zh-cn', 'ja', 'ru', 'ko', 'te', 'kn', 'ml', 'mr', 'gu', 'pa', 'bn']

def background_pre_translate(text, filename, source_lang):
    """
    Translates the book and its summary into all supported languages in a background thread
    to ensure instantaneous language switching for the user.
    """
    if not text: return
    
    path = os.path.join(UPLOAD_FOLDER, filename)
    log_error(f"[Pre-Translate] Starting background engine for '{filename}' from source: {source_lang}")
    
    # 1. Generate base summary if not exists
    base_summary = ""
    try:
        # We use a smaller portion for base summary to be fast
        summary_text = text
        if len(text) > 15000:
            summary_text = text[:10000] + "\n...\n" + text[-5000:]
        base_summary = summarize_text(f"Summarize the following text in exactly 50 words: {summary_text}")
    except Exception as se:
        log_error(f"[Pre-Translate] Base summary generation failed: {se}")

    for lang in PRE_TRANSLATE_LANGS:
        if lang == source_lang.lower(): continue
        
        try:
            # A. Process Book Translation
            trans_cache_path = get_translation_cache_path(path, lang)
            if not trans_cache_path or not os.path.exists(trans_cache_path):
                log_error(f"[Pre-Translate] Translating book '{filename}' to {lang}...")
                translated_text = translate_with_markers(text, lang, source_lang=source_lang)
                if trans_cache_path:
                    with open(trans_cache_path, "w", encoding="utf-8") as f:
                        json.dump({"translated_text": translated_text}, f)
            
            # B. Process Summary Translation
            sum_cache_filename = f"{hashlib.md5(filename.encode()).hexdigest()}_{lang}.json"
            sum_cache_path = os.path.join(SUMMARY_CACHE_DIR, sum_cache_filename)
            
            if base_summary and not os.path.exists(sum_cache_path):
                log_error(f"[Pre-Translate] Translating summary for '{filename}' to {lang}...")
                translated_summary = GoogleTranslator(source='auto', target=lang).translate(base_summary)
                with open(sum_cache_path, "w", encoding="utf-8") as f:
                    json.dump({"summary": translated_summary}, f)
            
            log_error(f"[Pre-Translate] Successfully cached book & summary for '{filename}' in {lang}")
            
            # Small delay to prevent API rate limiting
            time.sleep(2.5) 
            
        except Exception as e:
            log_error(f"[Pre-Translate] Failed for '{filename}' to {lang}: {e}")

@app.route("/api/read/<path:filename>")
def read_book(filename):
    path = os.path.join(UPLOAD_FOLDER, filename)
    lang = request.args.get('lang', 'en')

    if not os.path.exists(path):
        return jsonify({"error": "File not found"}), 404

    original_text, is_image, is_video = extract_text_from_file(path, filename)
    text = original_text
    
    # Language code to name mapping
    LANG_MAP = {
        'en': 'English', 'ta': 'Tamil', 'hi': 'Hindi', 'fr': 'French',
        'es': 'Spanish', 'de': 'German', 'zh-cn': 'Chinese', 'zh-tw': 'Chinese',
        'ja': 'Japanese', 'ru': 'Russian', 'ko': 'Korean', 'te': 'Telugu',
        'kn': 'Kannada', 'ml': 'Malayalam', 'mr': 'Marathi', 'gu': 'Gujarati',
        'pa': 'Punjabi', 'bn': 'Bengali'
    }

    detect_lang = 'en'
    detect_lang_name = 'English'
    
    if text:
        try:
            # Try to get text from Page 6 as requested for more representative sample
            # Page 1 often has titles/mixed text
            test_text = text
            if "--- Page 6 ---" in text:
                test_text = text.split("--- Page 6 ---", 1)[1][:2000]
            elif "--- Page 3 ---" in text:
                test_text = text.split("--- Page 3 ---", 1)[1][:2000]
            else:
                test_text = text[:2000]
                
            if test_text.strip():
                detect_lang = detect(test_text)
                detect_lang_name = LANG_MAP.get(detect_lang.lower(), detect_lang.upper())
                log_error(f"Detected language for {filename}: {detect_lang} ({detect_lang_name})")
        except Exception as e:
            log_error(f"Language detection error: {e}")

    # Check Translation Cache
    trans_cache_path = get_translation_cache_path(path, lang)
    if text and lang != detect_lang:
        if trans_cache_path and os.path.exists(trans_cache_path):
            log_error(f"Using cached translation for {filename} ({lang})")
            with open(trans_cache_path, "r", encoding="utf-8") as f:
                text = json.load(f).get('translated_text')
        else:
            # If user explicitly asked for a language that isn't the original
            # Pass the detected language as source for higher accuracy
            log_error(f"NO CACHE FOUND. Translating {filename} ({lang})...")
            text = translate_with_markers(text, lang, source_lang=detect_lang)
            # Save to translation cache
            if trans_cache_path:
                with open(trans_cache_path, "w", encoding="utf-8") as f:
                    json.dump({"translated_text": text}, f)

    is_pdf = filename.lower().endswith(".pdf")
    is_office = filename.lower().endswith((".docx", ".pptx", ".doc", ".ppt"))
    is_txt = filename.lower().endswith(".txt")
    
    preview_filename = filename
    office_html = ""
    if is_office:
        from utils.office_to_pdf import convert_to_pdf
        # Fast Path: Check if PDF already exists in cache
        base, _ = os.path.splitext(path)
        cached_pdf = base + ".pdf"
        
        if os.path.exists(cached_pdf):
            preview_filename = os.path.basename(cached_pdf)
            is_pdf = True
            is_office = False
        else:
            # First open: Use Mammoth/PPTX HTML for an instant experience
            office_html = generate_html_preview(path, filename, lang=lang)
            
            # Kick off the high-fidelity conversion in the background
            # We don't join/wait for it here, so the API returns immediately.
            def background_conversion(p, f):
                try:
                    convert_to_pdf(p) 
                    log_error(f"Background conversion finished for {f}")
                except Exception as bge:
                    log_error(f"Background conversion failed for {f}: {bge}")
            
            threading.Thread(target=background_conversion, args=(path, filename), daemon=True).start()
    
    elif is_txt:
        office_html = generate_html_preview(path, filename, lang=lang)

    # --- BACKGROUND PRE-TRANSLATE TRIGGER ---
    # Kick off the translation engine for ALL file types after the reader returns (asynchronously)
    threading.Thread(target=background_pre_translate, args=(original_text, filename, detect_lang), daemon=True).start()

    # Detect images (PDF or Docx) - Check original extension
    has_images = False
    orig_ext = filename.lower()
    if orig_ext.endswith(".docx"):
        try:
            from docx import Document
            doc = Document(path)
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    has_images = True
                    break
        except Exception as e:
            log_error(f"Docx image detection error for {filename}: {e}")
    elif orig_ext.endswith(".pdf") or (is_pdf and not orig_ext.endswith(".docx")):
        # If it was originally a PDF OR it was converted but we have the PDF path?
        # Actually, ocr.count_images(path) only works on PDF files.
        # If it was converted, we should check the converted PDF!
        check_path = path
        if is_pdf and orig_ext.endswith(".docx"):
            base, _ = os.path.splitext(path)
            conv_pdf = base + ".pdf"
            if os.path.exists(conv_pdf):
                check_path = conv_pdf
        
        has_images = ocr.count_images(check_path)

    # Get total page count
    num_pages = get_page_count(path, filename)

    return jsonify({
        "filename": filename,
        "preview_filename": preview_filename,
        "text": text,
        "original_text": original_text,
        "pages": int(num_pages) if str(num_pages).isdigit() else 1,
        "is_image": is_image,
        "is_pdf": is_pdf,
        "is_video": is_video,
        "is_office": is_office,
        "office_html": office_html,
        "lang": lang,
        "detected_lang": detect_lang,
        "detected_lang_name": detect_lang_name,
        "has_images": has_images
    })

@app.route("/api/extract_images", methods=["POST"])
def extract_images():
    filename = request.form.get("filename", "").strip()
    if not filename: return jsonify({"error": "No file specified"}), 400
    
    path = os.path.join(UPLOAD_FOLDER, filename)
    if not os.path.exists(path): return jsonify({"error": "File not found"}), 404
    
    # Check if it's an office file that was converted to PDF
    is_pdf = filename.lower().endswith(".pdf")
    if not is_pdf:
        from utils.office_to_pdf import convert_to_pdf
        path = convert_to_pdf(path) or path
        is_pdf = path.lower().endswith(".pdf")

    if is_pdf:
        output_dir = os.path.join("static", "extracted_images")
        results = ocr.extract_images_from_pdf(path, output_dir)
        return jsonify({"images": results})
    
    # Fallback for .docx if PDF conversion failed
    if filename.lower().endswith(".docx"):
        try:
            from docx import Document
            doc = Document(path)
            output_dir = os.path.join("static", "extracted_images")
            os.makedirs(output_dir, exist_ok=True)
            results = []
            img_count = 0
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    img_count += 1
                    img_name = f"{filename}_{img_count}.png"
                    img_path = os.path.join(output_dir, img_name)
                    with open(img_path, "wb") as f:
                        f.write(rel.target_part.blob)
                    results.append({
                        "url": f"/static/extracted_images/{img_name}",
                        "explanation": f"Image {img_count} from Word document."
                    })
            return jsonify({"images": results})
        except Exception as e:
            log_error(f"Docx image extraction error: {e}")
    
    return jsonify({"images": [], "message": "Image extraction only supported for PDF and Office documents."})

@app.route("/api/prepare_translation", methods=["POST"])
def prepare_translation():
    filename = request.form.get("filename")
    lang = request.form.get("lang", "en")
    
    if not filename:
        return jsonify({"error": "No filename"}), 400
        
    path = os.path.join(UPLOAD_FOLDER, filename)
    if not os.path.exists(path):
        return jsonify({"error": "File not found"}), 404
        
    # Trigger extraction (uses cache if available)
    text, is_image, is_video = extract_text_from_file(path, filename)
    
    # Detect language
    detect_lang = 'en'
    if text:
        try:
            test_text = text[:2000]
            if "--- Page 6 ---" in text:
                test_text = text.split("--- Page 6 ---", 1)[1][:2000]
            detect_lang = detect(test_text)
        except:
            pass
            
    # Check if translation is already cached
    trans_cache_path = get_translation_cache_path(path, lang)
    if text and lang != detect_lang:
        if trans_cache_path and os.path.exists(trans_cache_path):
            return jsonify({"status": "ready", "message": "Translation already cached."})
        else:
            # Trigger translation and cache it
            log_error(f"Preparing translation for {filename} to {lang}")
            translated_text = translate_with_markers(text, lang, source_lang=detect_lang)
            if trans_cache_path:
                with open(trans_cache_path, "w", encoding="utf-8") as f:
                    json.dump({"translated_text": translated_text}, f)
            return jsonify({"status": "completed", "message": "Translation prepared and cached."})
            
    return jsonify({"status": "error", "message": "Translation not needed or failed."})


@app.route("/api/summarize_file", methods=["POST"])
@app.route("/api/summarize_file/<path:filename>", methods=["POST"])
def summarize_file(filename=None):
    if not filename:
        filename = request.form.get("filename")
        
    path = os.path.join(UPLOAD_FOLDER, filename) if filename else None
    
    # Check if custom text was provided (selective summarization)
    custom_text = request.form.get("text")
    
    if custom_text:
        text = custom_text
    elif path and os.path.exists(path):
        text, is_image, is_video = extract_text_from_file(path, filename)
        # Truncate for summarization to avoid exceeding token limits
        if text and len(text) > 15000:
            text = text[:10000] + "\n...\n" + text[-5000:]
    else:
        return jsonify({"error": "File not found"}), 404

    lang = request.form.get("lang", "en")
    
    # Check Summary Cache First
    sum_cache_filename = f"{hashlib.md5(filename.encode()).hexdigest()}_{lang}.json"
    sum_cache_path = os.path.join(SUMMARY_CACHE_DIR, sum_cache_filename)
    
    if os.path.exists(sum_cache_path):
        try:
            with open(sum_cache_path, "r", encoding="utf-8") as f:
                cached_data = json.load(f)
                return jsonify({"summary": cached_data.get("summary")})
        except: pass

    try:
        summary = summarize_text(f"Summarize the following text in exactly 50 words: {text}")
        if lang != 'en' and summary:
            summary = GoogleTranslator(source='auto', target=lang).translate(summary)
            # Cache the result for future
            with open(sum_cache_path, "w", encoding="utf-8") as f:
                json.dump({"summary": summary}, f)
    except Exception as e:
        log_error(f"Summarization route error: {e}")
        summary = f"Error during summarization: {str(e)}"

    if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
        return jsonify({"summary": summary})

    return render_template(
        "summary.html",
        filename=filename,
        summary=summary
    )


@app.route("/uploads/<path:filename>")
def uploaded_file(filename):
    response = send_from_directory(UPLOAD_FOLDER, filename, as_attachment=False)
    # Add headers to help browser render content correctly in iframes/embeds
    if filename.lower().endswith(".pdf"):
        response.headers['Content-Type'] = 'application/pdf'
    response.headers['X-Content-Type-Options'] = 'nosniff'
    # Optional: ensure no cross-origin blockers
    response.headers['Access-Control-Allow-Origin'] = '*'
    return response


@app.route("/api/ask", methods=["POST"])
def ask():
    question = request.form.get("question", "").strip()
    context = request.form.get("context", "").strip()
    filename = request.form.get("filename", "").strip()

    lang = request.form.get("lang", "en")

    log_error(f"Ask AI triggered: {question} for {filename} (lang: {lang})")
    
    # Truncate context to avoid token limits (keep start and end for overall context)
    if context and len(context) > 20000:
        context = context[:10000] + "\n[... Content Truncated ...]\n" + context[-10000:]

    if not context:
        answer = "No readable text available for this file."
    elif not question:
        answer = "Please enter a question."
    else:
        try:
            # 1. Try Gemini Assistant IF API key looks valid
            from utils.gemini_assistant import GeminiAssistant
            api_key = os.getenv("GEMINI_API_KEY")
            
            answer = None
            if api_key and "YOUR_GEMINI_API_KEY" not in api_key and api_key.strip():
                answer = GeminiAssistant.ask(question, context)
            
            # 2. If Gemini is not configured or failed, fallback to local models
            if not answer:
                log_error(f"Gemini not available. Falling back to local model for: {question}")
                answer = answer_question(question, context)
            
            if lang != 'en' and answer:
                answer = GoogleTranslator(source='auto', target=lang).translate(answer)
        except Exception as e:
            log_error(f"Error in overall ask logic: {e}")
            answer = f"Error processing question: {str(e)}"

    if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
        return jsonify({"answer": answer})

    return f"""
    <h2>Answer</h2>
    <p><b>Question:</b> {question}</p>
    <p><b>Answer:</b> {answer}</p>
    <br>
    <a href="/read/{filename}">Back</a>
    """
@app.route("/api/speak", methods=["POST"])
def speak():
    text = request.form.get("text", "").strip()
    filename = request.form.get("filename", "").strip()
    lang = request.form.get("lang", "en").strip()

    rate = request.form.get("rate", "+0%").strip()
    gender = request.form.get("gender", "f").strip()

    if not text:
        return jsonify({"error": "No text available to convert into audio."})
    else:
        # Respect the expressive parameter from request, failing back to automatic check
        expressive_req = request.form.get("expressive", "true").lower() == "true"
        is_expressive = expressive_req and len(text.split()) > 3
        
        voice = request.form.get("voice")
        audio_file, vtt_file = text_to_speech(text, lang, rate=rate, gender=gender, expressive=is_expressive, voice=voice)
        
        if not audio_file:
             return jsonify({"error": vtt_file or "Speech generation failed"})
            
        return jsonify({
            "audio_url": f"/static/{audio_file}",
            "vtt_url": f"/static/{vtt_file}",
            "message": "Audio generated successfully"
        })



@app.route("/api/explain", methods=["POST"])
def explain():
    text = request.form.get("text", "").strip()
    lang = request.form.get("lang", "en").strip()
    filename = request.form.get("filename", "unknown")

    if not text:
        return jsonify({"error": "No text selected to explain."}), 400

    prompt = f"""
    Analyze the following paragraph and generate a clear, simple explanation.
    Target Language: {lang}
    
    Follow these rules strictly:
    1. Give a SIMPLE explanation:
       - Explain in very easy words
       - Make it understandable for a beginner
       - Avoid complex terms
    
    2. Give a REAL-LIFE example:
       - Relate the concept to everyday life
       - Keep it short and clear
    
    3. Extract KEY POINTS:
       - Provide 3 to 5 bullet points
       - Only include important ideas
    
    4. Keep the output clean and structured like this:
    
    Explanation:
    <simple explanation>
    
    Example:
    <real-life example>
    
    Key Points:
    - Point 1
    - Point 2
    - Point 3
    
    5. Do not add extra information
    6. Do not repeat the paragraph
    7. Keep the response concise and readable. Display professionally.
    
    Paragraph: 
    {text}
    """

    try:
        from utils.gemini_assistant import GeminiAssistant
        GeminiAssistant.configure()
        answer = GeminiAssistant.ask(prompt)
        
        # If Gemini is missing or failed, provide a much smarter dynamic heuristic fallback
        if not answer or "AI Error" in answer or "Gemini returned an empty response" in str(answer):
            from utils.qa import generate_explanation
            # Use local BART summarizer to find meaning if API key is missing
            try:
                summary = generate_explanation(text[:500], text)
                if summary and len(summary) > 20: 
                    answer = f"Explanation:\n{summary}\n\nExample:\nIt works just like a simplified version of {text.split()[0]} that you might use in daily life.\n\nKey Points:\n- Concepts are explained clearly.\n- Practical application is shown.\n- Core ideas are summarized."
                else:
                    # Better baseline if BART also fails
                    first_word = text.split()[0] if text else "context"
                    answer = f"Explanation:\nThis text discusses the core concept of '{first_word}'. It explores various implications of this subject.\n\nExample:\nImagine this being used in a real-world scenario where you need a simple tool.\n\nKey Points:\n- Fundamental ideas\n- Practical context\n- General understanding"
            except:
                answer = "Explanation:\nUnable to generate local explanation. Please check your Gemini API key in the .env file."
            
        return jsonify({"explanation": answer})
    except Exception as e:
        log_error(f"Explain API error: {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/meaning", methods=["POST"])
def meaning():
    word = request.form.get("word", "").strip()
    filename = request.form.get("filename", "unknown")
    context = request.form.get("context", "")
    lang = request.form.get("lang", "en").strip()

    log_error(f"Meaning lookup triggered for word: '{word}' in file: {filename} (lang: {lang})")
    if not word:
        result = "Please enter a word."
    else:
        try:
            log_error(f"Calling get_meaning for '{word}'...")
            result = get_meaning(word, context=context, target_lang=lang)
            log_error(f"get_meaning returned successfully for '{word}'")
        except Exception as e:
            log_error(f"CRITICAL Error in meaning route: {e}")
            result = f"Error looking up word: {str(e)}"

    if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
        log_error(f"Returning JSON meaning for '{word}'")
        return jsonify({"meaning": result})

    return f"""
    <h2>Word Meaning</h2>
    <p><b>{word}</b> : {result}</p>
    <br>
    <a href="/read/{filename}">Back</a>
    """


@app.route("/api/generate_quiz", methods=["POST"])
def generate_quiz():
    filename = request.form.get("filename", "").strip()
    lang = request.form.get("lang", "en").strip()
    if not filename:
        return jsonify({"error": "No file specified."})
    
    # Get text content (use the robust internal function)
    path = os.path.join(UPLOAD_FOLDER, filename)
    
    if not os.path.exists(path):
         return jsonify({"error": f"Book file not found: {filename}"})

    try:
        text, _, _ = extract_text_from_file(path, filename)
    except Exception as e:
        log_error(f"Internal extraction error: {e}")
        text = ""

    if not text or len(text) < 100:
        return jsonify({"error": "Could not extract enough text from book to generate a quiz."})

    # Prepare prompt for Gemini
    sample_text = text[:10000] + "\n...\n" + text[len(text)//2 : len(text)//2 + 10000]
    
    prompt = f"""
    Based ONLY on the following book content, generate 10 high-quality Multiple Choice Questions (MCQs).
    Language: {lang}
    
    Requirements:
    1. Each question must have exactly 4 options.
    2. Mark the correct answer index (0, 1, 2, or 3).
    3. Provide a brief explanation for why the answer is correct.
    4. Ensure questions are challenging but fair, covering key plot points or facts.
    5. Return ONLY a valid JSON array of objects.
    
    JSON Format:
    [
      {{
        "question": "What is the main theme of the first chapter?",
        "options": ["Option A", "Option B", "Option C", "Option D"],
        "answer": 0,
        "explanation": "Chapter 1 introduces... hence Option A."
      }},
      ...
    ]
    
    Book Content (Sample):
    {sample_text}
    """
    
    try:
        from utils.gemini_assistant import GeminiAssistant
        has_gemini = GeminiAssistant.configure()
        
        if has_gemini:
            response = GeminiAssistant.ask(prompt)
            if response and "AI Error" not in response and "Gemini returned an empty response" not in response:
                # Clean response string to ensure it's pure JSON
                import json
                clean_json = response.strip()
                if "```json" in clean_json:
                    clean_json = clean_json.split("```json")[1].split("```")[0].strip()
                elif "```" in clean_json:
                    clean_json = clean_json.split("```")[1].split("```")[0].strip()
                    
                try:
                    questions = json.loads(clean_json)
                    return jsonify({"questions": questions})
                except json.JSONDecodeError:
                    log_error(f"Failed to parse Gemini JSON: {clean_json}")
        
        # If Gemini failed or was not configured, use the "Free Service" (Heuristic Fallback)
        log_error("Using Free Basic Service for quiz generation (Heuristic fallback)")
        free_questions = []
        import re
        
        # Clean text of internal markers for professional display (Harden regex)
        clean_text = re.sub(r'(?i)\s*--- (?:Page|Slide) \d+ ---\s*', ' ', text)
        
        # Heuristic: find sentences that denote "importance" or "significance"
        impact_keywords = ['important', 'main', 'key', 'impact', 'result', 'benefit', 'consequence', 'goal', 'feature', 'significant', 'primary', 'enables', 'provides', 'represents']
        
        sentences = re.split(r'[.!?]\s+', clean_text)
        high_value_sentences = []
        normal_sentences = []
        
        for s in sentences:
            s_stripped = s.strip()
            if not s_stripped or len(s_stripped.split()) < 12: continue
            
            # Check for impact keywords
            score = sum(1 for kw in impact_keywords if re.search(r'\b' + kw + r'\b', s_stripped, re.IGNORECASE))
            if score > 0:
                high_value_sentences.append((s_stripped, score))
            else:
                normal_sentences.append(s_stripped)
        
        # Sort high-value by impact score
        high_value_sentences.sort(key=lambda x: x[1], reverse=True)
        
        # Combine lists to ensure we have enough for 10 questions
        # Prioritize high-value, then normal sentences
        candidate_sentences = [s for s, score in high_value_sentences] + normal_sentences
        import random
        
        # We need a pool of sentences for distractors
        distractor_pool = [s.strip() for s in sentences if len(s.strip().split()) > 8]
        
        for s in candidate_sentences:
            if len(free_questions) >= 10: break
            
            # Identify a potential subject
            verbs = r'\b(is|are|was|were|means|refers to|enables|provides|represents|impacts|allows|helps)\b'
            match = re.search(verbs, s, re.IGNORECASE)
            
            if match or len(free_questions) >= 5:
                # Subject Extraction
                if match:
                    subject_part = s[:match.start()].strip()
                    words = subject_part.split()
                    subject = f"'{subject_part}'" if 0 < len(words) <= 6 else "the topic detailed"
                else:
                    words = s.split()
                    subject = f"'{' '.join(words[:4])}...'" if len(words) > 4 else "the subject mentioned"
                
                # Correct Option
                correct_opt = s
                if len(correct_opt) > 160: correct_opt = correct_opt[:157] + "..."
                
                # Dynamic Distractors from other parts of the text
                others = [d for d in distractor_pool if d != s and len(d.split()) > 10]
                if len(others) >= 3:
                    chosen_distractors = random.sample(others, 3)
                else:
                    # Fallback distractors if text is too short
                    chosen_distractors = [
                        "This detail is not explicitly mentioned in the primary documentation.",
                        "The information provided suggests a contrary interpretation or legacy view.",
                        "Insufficient contextual data is available to support this specific observation."
                    ]
                
                # Clean and Truncate Distractors
                final_distractors = []
                for d in chosen_distractors:
                    if len(d) > 160: d = d[:157] + "..."
                    final_distractors.append(d)
                
                # Assemble and Shuffle
                options = [correct_opt] + final_distractors
                random.shuffle(options)
                ans_idx = options.index(correct_opt)
                
                free_questions.append({
                    "question": f"Based on the provided documentation, what is a primary significance or impact of {subject}?",
                    "options": options,
                    "answer": ans_idx,
                    "explanation": f"The source text emphasizes: \"{s}\""
                })
        
        if not free_questions:
            return jsonify({"error": "Book content is too short to generate even a basic quiz. Please try a different book."})
            
        # Translate to target language if not English
        if lang != 'en':
            try:
                translator = GoogleTranslator(source='auto', target=lang)
                for q in free_questions:
                    q['question'] = translator.translate(q['question'])
                    q['options'] = [translator.translate(opt) for opt in q['options']]
                    q['explanation'] = translator.translate(q['explanation'])
            except Exception as te:
                log_error(f"Fallback translation error: {te}")

        return jsonify({
            "questions": free_questions,
            "message": f"AI Service currently unavailable. Using basic free generation in {lang} instead."
        })
    except Exception as e:
        log_error(f"Quiz generation error: {e}")
        return jsonify({"error": f"Failed to generate quiz: {str(e)}"})

@app.route("/api/save_quiz_score", methods=["POST"])
def save_quiz_score():
    filename = request.form.get("filename", "").strip()
    score = request.form.get("score", "0")
    total = request.form.get("total", "10")
    
    try:
        conn = get_db_connection()
        if not conn:
            return jsonify({"success": False, "error": "Database unavailable"}), 503
            
        cursor = conn.cursor()
        cursor.execute(
            "INSERT INTO quiz_history (filename, score, total) VALUES (%s, %s, %s)",
            (filename, score, total)
        )
        conn.commit()
        conn.close()
        return jsonify({"success": True})
    except Exception as e:
        log_error(f"Save quiz score error: {e}")
        return jsonify({"error": str(e)})

@app.route("/api/recommendations")
def get_recommendations():
    # Curated recommendations for different categories across multiple languages
    recommendations = {
        "Story": {
            "English": [
                {"title": "The Alchemist", "author": "Paulo Coelho", "desc": "A journey of self-discovery and following one's destiny."},
                {"title": "Alice in Wonderland", "author": "Lewis Carroll", "desc": "A whimsical journey through a fantastical world."},
                {"title": "The Little Prince", "author": "Antoine de Saint-Exupéry", "desc": "A profound tale about life and human nature."}
            ],
            "Tamil": [
                {"title": "Panchatantra Stories", "author": "Vishnu Sharma", "desc": "Ancient Indian animal fables with moral lessons."},
                {"title": "Siruvar Kathaigal", "author": "Various", "desc": "Captivating short stories for cognitive and moral growth."},
                {"title": "Tenali Raman Tales", "author": "Traditional", "desc": "Witty and humorous stories of the legendary courtier."}
            ],
            "Hindi": [
                {"title": "Panchatantra", "author": "Vishnu Sharma", "desc": "Educational and moral fables for all ages."},
                {"title": "Akbar-Birbal", "author": "Traditional", "desc": "Stories of wisdom and quick wit."},
                {"title": "Vikram-Betal", "author": "Traditional", "desc": "Mystical tales of King Vikram and the ghost Betal."}
            ]
        },
        "Novels": {
            "English": [
                {"title": "1984", "author": "George Orwell", "desc": "A dystopian look at a totalitarian future society."},
                {"title": "Pride and Prejudice", "author": "Jane Austen", "desc": "A classic romance exploring manners and social standing."},
                {"title": "The Great Gatsby", "author": "F. Scott Fitzgerald", "desc": "A tale of the American Dream in the Roaring Twenties."}
            ],
            "Tamil": [
                {"title": "Ponniyin Selvan", "author": "Kalki", "desc": "Epic historical novel about the Chola Empire."},
                {"title": "Sivagamiyin Sabatham", "author": "Kalki", "desc": "A historical masterpiece set in the Pallava kingdom."},
                {"title": "Parthiban Kanavu", "author": "Kalki", "desc": "A gripping tale of a king seeking independence."}
            ],
            "Hindi": [
                {"title": "Godaan", "author": "Premchand", "desc": "A powerful commentary on socio-economic struggles in rural India."},
                {"title": "Gaban", "author": "Premchand", "desc": "An exploration of social pressure and personal integrity."},
                {"title": "Gunahon Ka Devta", "author": "Dharamvir Bharati", "desc": "An iconic emotional and tragic love story."}
            ]
        },
        "Inspirational": {
            "English": [
                {"title": "Wings of Fire", "author": "A.P.J. Abdul Kalam", "desc": "The autobiography of the Missile Man of India."},
                {"title": "Think and Grow Rich", "author": "Napoleon Hill", "desc": "Principles for achieving success and wealth."},
                {"title": "Man's Search for Meaning", "author": "Viktor Frankl", "desc": "Finding purpose even in the darkest circumstances."}
            ],
            "Tamil": [
                {"title": "Agni Siragugal", "author": "A.P.J. Abdul Kalam", "desc": "Tamil version of the inspiring Wings of Fire."},
                {"title": "Arthamulla Indhu Madham", "author": "Kannadasan", "desc": "Philosophical insights into life and spirituality."},
                {"title": "Thirukkural", "author": "Thiruvalluvar", "desc": "Ancient Tamil ethics and wisdom in couplets."}
            ],
            "Hindi": [
                {"title": "Agni Ki Udaan", "author": "A.P.J. Abdul Kalam", "desc": "Inspiring life journey of the former President."},
                {"title": "Jeet Aapki", "author": "Shiv Khera", "desc": "Practical steps to achieve positive attitude and success."},
                {"title": "Madhushala", "author": "Harivansh Rai Bachchan", "desc": "Philosophical poetry about the essence of life's path."}
            ]
        },
        "Comedy": {
            "English": [
                {"title": "The Hitchhiker's Guide to the Galaxy", "author": "Douglas Adams", "desc": "A hilarious cosmic journey across the universe."},
                {"title": "Catch-22", "author": "Joseph Heller", "desc": "Satirical look at the absurdity of war and bureaucracy."},
                {"title": "Three Men in a Boat", "author": "Jerome K. Jerome", "desc": "A humorous account of a boating holiday on the Thames."}
            ],
            "Tamil": [
                {"title": "Washingtonil Thirumanam", "author": "Savvi", "desc": "A classic humorous look at NRI life experiences."},
                {"title": "Thuppariyum Sambu", "author": "Devan", "desc": "The comical adventures of an accidental detective."},
                {"title": "Mudra Rakshasam", "author": "Cho Ramaswamy", "desc": "Sharp political satire and wit."}
            ],
            "Hindi": [
                {"title": "Raag Darbari", "author": "Shrilal Shukla", "desc": "Satirical masterpiece on rural politics and life."},
                {"title": "Mungerilal Ke Haseen Sapne", "author": "Manohar Shyam Joshi", "desc": "Funny tales of a common man's daydreams."},
                {"title": "Chitchor", "author": "Various", "desc": "Lighthearted stories of romance and misunderstandings."}
            ]
        }
    }
    return jsonify(recommendations)
    
@app.route("/api/predict_questions", methods=["POST"])
def predict_questions():
    filename = request.form.get("filename", "").strip()
    lang = request.form.get("lang", "en").strip()
    
    if not filename:
        return jsonify({"error": "No file specified."}), 400
        
    path = os.path.join(UPLOAD_FOLDER, filename)
    if not os.path.exists(path):
        return jsonify({"error": "File not found."}), 404
        
    try:
        text, _, _ = extract_text_from_file(path, filename)
    except Exception as e:
        log_error(f"Extraction error for prediction: {e}")
        return jsonify({"error": "Could not extract text."}), 500
        
    if not text or len(text) < 100:
        return jsonify({"error": "Book content too short for prediction."}), 400

    # Prepare prompt for Gemini
    sample_text = text[:15000] # Use a bit more for prediction if possible
    
    prompt = f"""
    Analyze the following book content and generate exactly 5 'Short Answer Questions' and 3 'Long Answer Questions' for an exam.
    Target Language: {lang}
    
    Requirements:
    1. Focus on the MOST important concepts, repeated themes, and key facts.
    2. Short Answer Questions: Each answer MUST be EXACTLY 4 lines long. Use line breaks ("\n") to separate lines.
    3. Long Answer Questions: Each answer MUST be EXACTLY 8 lines long. Use line breaks ("\n") to separate lines.
    4. For EACH answer, use **bold text** (e.g., **Key Term**) for the most important keywords or phrases to highlight them.
    5. For EACH question, provide:
       - The question text
       - A clear, structured answer/explanation (adhering strictly to the line count rule)
       - Difficulty level: "Easy", "Medium", or "Hard"
       - A boolean "is_important" (set to true for the top 2-3 most critical overall questions)
    6. Return ONLY a valid JSON object in this format:
    {{
      "short_questions": [
        {{ "question": "...", "answer": "...", "difficulty": "...", "is_important": true }},
        ...
      ],
      "long_questions": [
        {{ "question": "...", "answer": "...", "difficulty": "...", "is_important": false }},
        ...
      ]
    }}
    
    Book Content (Sample):
    {sample_text}
    """

    try:
        from utils.gemini_assistant import GeminiAssistant
        has_gemini = GeminiAssistant.configure()
        
        if has_gemini:
            response = GeminiAssistant.ask(prompt)
            if response and "AI Error" not in response and "Gemini returned an empty response" not in response:
                clean_json = response.strip()
                if "```json" in clean_json:
                    clean_json = clean_json.split("```json")[1].split("```")[0].strip()
                elif "```" in clean_json:
                    clean_json = clean_json.split("```")[1].split("```")[0].strip()
                
                try:
                    data = json.loads(clean_json)
                    return jsonify(data)
                except Exception as je:
                    log_error(f"JSON Parse Error in prediction: {je}")
    except Exception as e:
        log_error(f"Gemini Prediction Error: {e}")

    # Fallback / Simulated Logic (Heuristic)
    log_error("Using heuristic fallback for Smart Prediction")
    import re
    import random
    
    clean_text = re.sub(r'(?i)\s*--- (?:Page|Slide) \d+ ---\s*', ' ', text)
    sentences = [s.strip() for s in re.split(r'[.!?]\s+', clean_text) if len(s.strip().split()) > 10]
    
    if not sentences:
        return jsonify({"error": "Insufficient content for prediction."}), 400

    # Find "important" sentences
    keywords = ['important', 'key', 'main', 'significance', 'result', 'defined', 'because', 'therefore', 'consequently']
    important_sentences = []
    for s in sentences:
        score = sum(1 for kw in keywords if kw in s.lower())
        if score > 0:
            important_sentences.append((s, score))
    
    important_sentences.sort(key=lambda x: x[1], reverse=True)
    candidates = [s for s, score in important_sentences[:15]]
    if not candidates: candidates = sentences[:10]
    
    short_qs = []
    for i, s in enumerate(random.sample(candidates, min(len(candidates), 5))):
        words = s.split()
        short_qs.append({
            "question": f"What is the significance of the concept mentioned as: '{' '.join(words[:6])}...'?",
            "answer": f"The book states: '{s}'.\nThis concept is crucial for understanding the primary subject matter.\nIt highlights key details that are fundamental to the chapter's progression.\nOverall, this point serves as a foundation for further analysis of the text.",
            "difficulty": random.choice(["Easy", "Medium"]),
            "is_important": i < 2
        })
        
    long_qs = []
    for i, s in enumerate(random.sample(candidates, min(len(candidates), 3))):
        words = s.split()
        long_qs.append({
            "question": f"Explain in detail the following point from the text: '{' '.join(words[:8])}...'",
            "answer": f"Based on the text: '{s}'.\nThis concept plays a vital role in the overall theme presented in this section.\nFurthermore, it provides a unique perspective on the historical or factual context.\nIt allows readers to connect various independent facts into a cohesive narrative structure.\nMoreover, the author emphasizes this specific idea to underline the central purpose of the study.\nBy examining this detail, one can derive a deeper meaning of the intended message.\nIt effectively encapsulates the main objectives while offering substantial evidence for its claims.\nUltimately, this analysis provides an 8-line comprehensive look as requested for long answers.",
            "difficulty": "Hard",
            "is_important": i == 0
        })
        
    # Translate to target language if not English
    if lang != 'en':
        try:
            translator = GoogleTranslator(source='auto', target=lang)
            for q in short_qs:
                q['question'] = translator.translate(q['question'])
                q['answer'] = translator.translate(q['answer'])
                q['difficulty'] = translator.translate(q['difficulty'])
            for q in long_qs:
                q['question'] = translator.translate(q['question'])
                q['answer'] = translator.translate(q['answer'])
                q['difficulty'] = translator.translate(q['difficulty'])
        except Exception as te:
            log_error(f"Prediction fallback translation error: {te}")

    return jsonify({
        "short_questions": short_qs,
        "long_questions": long_qs,
        "message": f"AI Service currently unavailable. Using basic prediction in {lang} instead."
    })

@app.route("/api/voices", methods=["GET"])
def get_voices():
    import asyncio
    import edge_tts
    try:
        # We use a helper to run the async call
        async def fetch_voices():
            return await edge_tts.list_voices()
        
        voices = asyncio.run(fetch_voices())
        return jsonify(voices)
    except Exception as e:
        log_error(f"Error fetching voices: {e}")
        return jsonify([])

if __name__ == "__main__":
    # Ensure tables exist (Simplified check)
    try:
        conn = get_db_connection()
        if conn:
            cursor = conn.cursor()
            cursor.execute("CREATE TABLE IF NOT EXISTS uploaded_books (id INT AUTO_INCREMENT PRIMARY KEY, filename VARCHAR(255), upload_time DATETIME, pages VARCHAR(50), last_opened DATETIME, last_closed DATETIME)")
            cursor.execute("CREATE TABLE IF NOT EXISTS quiz_history (id INT AUTO_INCREMENT PRIMARY KEY, filename VARCHAR(255), score INT, total INT, timestamp DATETIME DEFAULT CURRENT_TIMESTAMP)")
            conn.commit()
            conn.close()
    except: pass
    
    app.run(debug=True, port=int(os.environ.get("PORT", 5000)))